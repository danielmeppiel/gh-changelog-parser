import os
import re
import requests
from io import BytesIO
from parser_web import Article
from ai_tools import summarize_text
from bs4 import BeautifulSoup
from pptx import Presentation

MAX_ARTICLE_TITLE_WORDS = os.environ.get('MAX_ARTICLE_TITLE_WORDS', 10)
MAX_ARTICLE_CONTENT_CHARS = os.environ.get('MAX_ARTICLE_CONTENT_CHARS', 215)

TITLE_TMPL_SLIDE_IDX = os.environ.get('TITLE_TMPL_SLIDE_IDX', 0)
AGENDA_TMPL_SLIDE_IDX = os.environ.get('AGENDA_TMPL_SLIDE_IDX', 1)
THX_TMPL_SLIDE_IDX = os.environ.get('AGENDA_TMPL_SLIDE_IDX', 9)

ARTICLE_W_IMAGE_SLIDE_IDX = os.environ.get('ARTICLE_W_IMAGE_SLIDE_IDX', 5)
ARTICLE_WO_IMAGE_SLIDE_IDX = os.environ.get('ARTICLE_WO_IMAGE_SLIDE_IDX', 3)

DEVEX_SLIDE_MASTER_IDX = os.environ.get('DEVEX_SLIDE_MASTER_IDX', 0)
SECURITY_SLIDE_MASTER_IDX = os.environ.get('SECURITY_SLIDE_MASTER_IDX', 2)
AI_SLIDE_MASTER_IDX = os.environ.get('AI_SLIDE_MASTER_IDX', 1)

CATEGORY_SLIDE_TMPL_IDX = os.environ.get('CATEGORY_SLIDE_TMPL_IDX', 2)

def parse_changelog(file_path):
    articles = []
    current_category = None
    with open(file_path, 'r') as file:
        for line in file:
            category_match = re.match(r'#+\s(.+)', line)
            if category_match:
                current_category = category_match.group(1)
            else:
                article_match = re.match(r'- \((.*?)\) \[(.*?)\]\((.*?)\)', line)
                if article_match:
                    published, title, link = article_match.groups()
                    article = Article(title=title, link=link, published=published, category=current_category)
                    articles.append(article)
    return articles

def fetch_article_content(article):
    response = requests.get(article.link)
    response.raise_for_status()
    soup = BeautifulSoup(response.content, 'html.parser')
    # Assuming the article content is within a <div> with class 'post-content'
    content_div = soup.find('section', class_='post__content')
    if content_div:
        paragraphs = content_div.find_all('p')
        article.summary = ' '.join([p.get_text(strip=True) for p in paragraphs])
    
    # Extract img elements from content_div, if any, and append to imgs
    img_elements = content_div.find_all('img')
    for img in img_elements:
        article.imgs.append(img)
    return article

def get_category_master_slide(prs, category):
    if category == "DevEx":
        return prs.slide_masters[DEVEX_SLIDE_MASTER_IDX]
    elif category == "Security":
        return prs.slide_masters[SECURITY_SLIDE_MASTER_IDX]
    elif category == "AI":
        return prs.slide_masters[AI_SLIDE_MASTER_IDX]
    else:
        return prs.slide_masters[DEVEX_SLIDE_MASTER_IDX]

def choose_article_layout(article, prs):
    slide_master = get_category_master_slide(prs, article.category)
    slide_idx = ARTICLE_W_IMAGE_SLIDE_IDX if article.imgs else ARTICLE_WO_IMAGE_SLIDE_IDX
    return slide_master.slide_layouts[slide_idx]

def choose_category_layout(prs, category):
    slide_master = get_category_master_slide(prs, category)
    return slide_master.slide_layouts[CATEGORY_SLIDE_TMPL_IDX]

def add_agenda_slide(prs):
    print("Adding agenda slide...")
    slide_layout = prs.slide_layouts[AGENDA_TMPL_SLIDE_IDX]
    agenda_slide = prs.slides.add_slide(slide_layout)
    # Take the first placeholder of the agenda slide and add lines "01" "02" and "03"
    agenda_slide.placeholders[11].text = "01\n02\n03"
    # Take the second placeholder and add lines "DevEx" "Security" and "AI"
    agenda_slide.placeholders[12].text = "DevEx\nSecurity\nAI"


def add_category_slide(prs, category):
    print("Adding category slide for:", category)
    category_slide_layout = choose_category_layout(prs, category)
    category_slide = prs.slides.add_slide(category_slide_layout)
    title_placeholder = category_slide.placeholders[0]
    title_placeholder.text = category
    return category_slide

def add_article_slide(prs, article):
    print("Fetching content for:", article.title)
    slide_layout = choose_article_layout(article, prs)
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.placeholders[30]
    date_placeholder = slide.placeholders[31]
    content_placeholder = slide.placeholders[26]

    title_run = title_placeholder.text_frame.paragraphs[0].add_run()
    if len(article.title.split()) > MAX_ARTICLE_TITLE_WORDS:
        print("Summarizing title...")
        title_run.text = summarize_text(article.title, MAX_ARTICLE_TITLE_WORDS, "words")
    else:
        title_run.text = article.title
    title_run.hyperlink.address = article.link
    date_placeholder.text = article.published

    print("Summarizing content...")
    content_placeholder.text = summarize_text(article.summary, MAX_ARTICLE_CONTENT_CHARS)

    if article.imgs:
        img = article.imgs[0]
        img_url = img.get('src')
        fetch_and_insert_image(slide, img_url)

def fetch_and_insert_image(slide, img_url):
    print("Fetching image:", img_url)
    img_response = requests.get(img_url)
    img_stream = BytesIO(img_response.content)
    slide.placeholders[39].insert_picture(img_stream)

def update_title_slide(prs):
    print("Updating title slide...")
    title_slide = prs.slides[TITLE_TMPL_SLIDE_IDX]
    title_placeholder = title_slide.placeholders[0]
    title_placeholder.text = "Product Roadmap"

def add_thank_you_slide(prs):
    print("Adding thank you slide...")
    slide_layout = prs.slide_layouts[THX_TMPL_SLIDE_IDX]
    prs.slides.add_slide(slide_layout)

def create_presentation(changelog_items, template_file, output_file):
    prs = Presentation(template_file)
    update_title_slide(prs)
    add_agenda_slide(prs)
    category_slides = {}

    for article in changelog_items:
        if article.category not in category_slides:
            category_slide = add_category_slide(prs, article.category)
            category_slides[article.category] = category_slide
        article = fetch_article_content(article)
        add_article_slide(prs, article)

    add_thank_you_slide(prs)
    prs.save(output_file)

def main():
    # Read changelog data
    changelog_items = parse_changelog('changelog.md')

    # Create a new PowerPoint presentation using a template
    create_presentation(changelog_items, 'template.pptx', 'Changelog_Presentation.pptx')

    print("Presentation saved: Changelog_Presentation.pptx")

if __name__ == '__main__':
    main()